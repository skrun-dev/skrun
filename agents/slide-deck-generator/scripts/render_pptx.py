#!/usr/bin/env python3
"""Render a structured slide deck to deck.pptx via python-pptx.

Input  (stdin JSON):  { slides: [{ layout, title, subtitle?, bullets?, speaker_notes? }],
                        brand_primary_color?, deck_title? }
Output (stdout JSON): { path, bytes, slide_count }
"""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError as e:
    print(json.dumps({
        "error": f"python-pptx not installed: {e}. Run: pip install -r requirements.txt"
    }))
    sys.exit(0)


def parse_hex_color(hex_str: str) -> RGBColor:
    """Parse '#RRGGBB' into an RGBColor; default to #1a73e8 on parse error."""
    s = hex_str.strip().lstrip("#")
    if len(s) != 6:
        s = "1a73e8"
    try:
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return RGBColor(0x1a, 0x73, 0xe8)


def add_accent_bar(slide, color: RGBColor) -> None:
    """Top-left accent bar — visual signature of the brand color."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.25),
        Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def style_title(text_frame, color: RGBColor) -> None:
    """Apply brand color to a title text frame."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = color
            run.font.size = Pt(36)
            run.font.bold = True


def add_title_slide(prs: Presentation, title: str, subtitle: str | None, color: RGBColor) -> None:
    layout = prs.slide_layouts[0]  # title layout
    slide = prs.slides.add_slide(layout)
    add_accent_bar(slide, color)
    slide.shapes.title.text = title
    style_title(slide.shapes.title.text_frame, color)
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    speaker_notes: str | None,
    color: RGBColor,
) -> None:
    layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(layout)
    add_accent_bar(slide, color)
    slide.shapes.title.text = title
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = color
            run.font.bold = True

    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    if body is None:
        return
    tf = body.text_frame
    safe_bullets = bullets if bullets else ["_(no points captured)_"]
    tf.text = safe_bullets[0]
    for b in safe_bullets[1:]:
        p = tf.add_paragraph()
        p.text = b

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes


def add_closing_slide(prs: Presentation, title: str, subtitle: str | None, color: RGBColor) -> None:
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    add_accent_bar(slide, color)
    if slide.shapes.title:
        slide.shapes.title.text = title
        style_title(slide.shapes.title.text_frame, color)
    if subtitle:
        # Add a sub-text box below the title.
        left, top, width, height = Inches(1), Inches(3.5), Inches(8), Inches(1)
        box = slide.shapes.add_textbox(left, top, width, height)
        box.text_frame.text = subtitle
        for para in box.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(20)


def main() -> None:
    raw = sys.stdin.read()
    try:
        args = json.loads(raw)
    except json.JSONDecodeError as e:
        print(json.dumps({"error": f"invalid stdin JSON: {e}"}))
        return

    slides = args.get("slides", [])
    if not isinstance(slides, list) or not slides:
        print(json.dumps({"error": "slides must be a non-empty array"}))
        return

    color = parse_hex_color(args.get("brand_primary_color", "#1a73e8"))
    deck_title = args.get("deck_title", "")

    output_dir = os.environ.get("SKRUN_OUTPUT_DIR")
    if not output_dir:
        print(json.dumps({"error": "SKRUN_OUTPUT_DIR is not set"}))
        return

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    out_path = Path(output_dir) / "deck.pptx"

    prs = Presentation()
    if deck_title:
        prs.core_properties.title = deck_title

    for s in slides:
        layout = s.get("layout", "content")
        title = s.get("title", "Untitled")
        subtitle = s.get("subtitle")
        bullets = s.get("bullets") or []
        notes = s.get("speaker_notes")

        if layout == "title":
            add_title_slide(prs, title, subtitle, color)
        elif layout == "closing":
            add_closing_slide(prs, title, subtitle, color)
        else:
            add_content_slide(prs, title, list(bullets), notes, color)

    prs.save(str(out_path))
    size = out_path.stat().st_size
    print(json.dumps({"path": str(out_path), "bytes": size, "slide_count": len(slides)}))


if __name__ == "__main__":
    main()
